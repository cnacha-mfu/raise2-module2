"""
check_slides.py — ตรวจไฟล์ .pptx ที่สร้างแล้ว ว่ามีอะไรผิดพลาดหรือไม่

วิธีใช้:
    python scripts/check_slides.py materials/week5/w5-slides.pptx ...
    python scripts/check_slides.py            (ตรวจทุกสัปดาห์)

ตรวจ 4 อย่าง
  1. รูปทรงล้นขอบสไลด์ (ล่าง / ขวา)
  2. ข้อความยาวเกินกล่องที่รองรับ (ประมาณจากจำนวนตัวอักษรและขนาดฟอนต์)
  3. ข้อความจัดกึ่งกลางในที่ที่ควรชิดซ้าย
  4. สไลด์ที่ไม่มีบันทึกผู้บรรยาย
"""

import glob
import sys
from pathlib import Path

# หน้าจอ terminal ของ Windows ไทยใช้ cp874 ซึ่งพิมพ์ emoji ไม่ได้และจะ crash
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except (AttributeError, ValueError):
        pass

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOL = Inches(0.04)          # เผื่อความคลาดเคลื่อนเล็กน้อย
FOOTER_TOP = Inches(6.9)    # แถบเลขหน้าอยู่ตรงนี้

NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def para_align(p):
    pPr = p._p.find(f"{NS}pPr")
    return pPr.get("algn") if pPr is not None else None


def approx_text_height(shape):
    """ประมาณความสูงที่ข้อความต้องใช้จริง เทียบกับความสูงของกล่อง"""
    tf = shape.text_frame
    total = Emu(0)
    width_in = Emu(shape.width).inches - 0.2
    if width_in <= 0:
        return Emu(0)
    for p in tf.paragraphs:
        text = p.text
        size = None
        for r in p.runs:
            if r.font.size:
                size = r.font.size
                break
        size = size or Pt(18)
        # ตัวอักษรไทยกว้างประมาณ 0.55 เท่าของขนาดฟอนต์
        char_w_in = (size.pt * 0.55) / 72
        per_line = max(1, int(width_in / char_w_in))
        lines = max(1, -(-len(text) // per_line))
        total += Inches(size.pt * 1.35 / 72) * lines
    return total


def check(path):
    prs = Presentation(str(path))
    problems = []

    for idx, slide in enumerate(prs.slides, start=1):
        # สไลด์คั่นบทมีพื้นหลังเต็มจอ · สไลด์แบบนี้จัดกึ่งกลางถูกต้องแล้ว
        is_section = any(
            sh.width and sh.width > Inches(13) and sh.height and sh.height > Inches(7)
            for sh in slide.shapes
        )
        # สไลด์ข้อความใหญ่ = มีกล่องข้อความเดียวขนาดใหญ่ตรงกลาง
        texts = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        is_big = (len(texts) == 2 and idx > 1
                  and any(Emu(sh.height).inches > 2.5 for sh in texts))

        for sh in slide.shapes:
            if sh.top is None or sh.height is None:
                continue
            bottom, right = sh.top + sh.height, sh.left + sh.width
            full_bleed = sh.width > Inches(13) and sh.height > Inches(7)

            # 1) ล้นขอบสไลด์
            if bottom > SLIDE_H + TOL:
                problems.append((idx, "ล้นขอบล่าง",
                                 f"{Emu(bottom).inches:.2f}\" (สไลด์สูง 7.50\")"))
            if right > SLIDE_W + TOL:
                problems.append((idx, "ล้นขอบขวา",
                                 f"{Emu(right).inches:.2f}\" (สไลด์กว้าง 13.33\")"))

            # 2) ทับแถบเลขหน้า (ยกเว้นพื้นหลังเต็มจอของสไลด์คั่นบท)
            if not full_bleed and sh.height > Inches(0.5) and bottom > FOOTER_TOP + TOL:
                problems.append((idx, "ทับแถบเลขหน้า",
                                 f"ล่างสุดอยู่ที่ {Emu(bottom).inches:.2f}\""))

            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue

            # 3) ข้อความล้นกล่อง
            need = approx_text_height(sh)
            if need > sh.height + Inches(0.2):
                problems.append((idx, "ข้อความล้นกล่อง",
                                 f"ต้องการ ~{Emu(need).inches:.2f}\" "
                                 f"แต่กล่องสูง {Emu(sh.height).inches:.2f}\" "
                                 f"| {sh.text_frame.text.strip()[:32]}"))

            # 4) จัดกึ่งกลางในที่ที่ควรชิดซ้าย
            #    ที่กึ่งกลางถูกต้องแล้ว: สไลด์คั่นบท · สไลด์ข้อความใหญ่ ·
            #    กล่องในรูปวาด · ลูกศรวนกลับของ cycle
            is_diagram_box = sh.width < Inches(6.2) and sh.height < Inches(1.4)
            is_cycle_bar = sh.height < Inches(0.45) and sh.width > Inches(8)
            if not (is_section or is_big or is_diagram_box or is_cycle_bar):
                for p in sh.text_frame.paragraphs:
                    if para_align(p) == "ctr" and len(p.text.strip()) > 3:
                        problems.append((idx, "จัดกึ่งกลางผิดที่",
                                         f"{p.text.strip()[:38]!r}"))
                        break

    # สไลด์ปก (หน้า 1) ไม่ต้องมีบันทึกผู้บรรยาย
    missing_notes = [i for i, s in enumerate(prs.slides, start=1)
                     if i > 1 and not s.notes_slide.notes_text_frame.text.strip()]
    return problems, missing_notes, len(prs.slides._sldIdLst)


def main():
    targets = sys.argv[1:] or sorted(glob.glob("materials/week*/w*-slides.pptx"))
    if not targets:
        print("ไม่พบไฟล์ .pptx"); sys.exit(1)

    grand = 0
    for t in targets:
        problems, missing, n = check(Path(t))
        grand += len(problems)
        name = Path(t).name
        if not problems and not missing:
            print(f"✅ {name}  ({n} สไลด์) — ไม่พบปัญหา")
            continue
        print(f"\n📄 {name}  ({n} สไลด์)")
        if problems:
            print(f"   ❌ พบ {len(problems)} ปัญหา")
            for page, kind, detail in problems:
                print(f"      หน้า {page:>3} | {kind:<18} | {detail}")
        if missing:
            print(f"   ⚠️  สไลด์ที่ไม่มีบันทึกผู้บรรยาย: {missing}")

    print(f"\n{'='*60}")
    print("✅ ทุกไฟล์ผ่าน" if grand == 0 else f"❌ รวมพบ {grand} ปัญหา")
    sys.exit(0 if grand == 0 else 1)


if __name__ == "__main__":
    main()
