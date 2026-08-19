"""
build_slides.py — แปลงต้นฉบับสไลด์ (.md) เป็นไฟล์ PowerPoint (.pptx)

วิธีใช้:
    pip install python-pptx
    python scripts/build_slides.py materials/week5/w5-slides.md

ผลลัพธ์: ไฟล์ .pptx ชื่อเดียวกัน ในโฟลเดอร์เดียวกัน

รูปแบบต้นฉบับ (ดู CLAUDE.md หัวข้อ 6.2):
    ---
    week: 6
    title: "..."
    subtitle: "..."
    date: "5 กันยายน 2026"
    ---

    ## ชื่อช่วง            <- สไลด์คั่นบท
    # หัวข้อสไลด์          <- สไลด์เนื้อหา
    - bullet
    !! ข้อความใหญ่กลางจอ
    ```code
    โครงสร้าง/คำสั่ง
    ```
    ![screenshot: คำอธิบาย](assets/w6-01.png)
    > note: บันทึกผู้บรรยาย
"""

import re
import sys
from pathlib import Path

# หน้าจอ terminal ของ Windows ไทยใช้ cp874 ซึ่งพิมพ์ emoji ไม่ได้และจะ crash
# บังคับให้ใช้ utf-8 เอง ผู้ใช้จะได้ไม่ต้องตั้ง PYTHONIOENCODING ก่อนรันทุกครั้ง
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except (AttributeError, ValueError):
        pass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ──────────────────────────── ตั้งค่าที่แก้ได้ ────────────────────────────

# ⚠️ ฟอนต์ไทย — "Leelawadee UI" มีในเครื่อง Windows ทุกเครื่อง
#    ถ้าติดตั้ง Sarabun แล้วและอยากใช้ ให้เปลี่ยนบรรทัดนี้บรรทัดเดียว
THAI_FONT = "Leelawadee UI"
MONO_FONT = "Consolas"

COLOR_TEXT     = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_ACCENT   = RGBColor(0x1F, 0x5F, 0xA8)
COLOR_MUTED    = RGBColor(0x6B, 0x6B, 0x6B)
COLOR_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CODE_BG  = RGBColor(0xF2, 0xF4, 0xF7)
COLOR_IMG_BG   = RGBColor(0xE4, 0xE7, 0xEB)

# สีสำหรับรูปวาด (diagram)
COLOR_BOX_BG   = RGBColor(0xE8, 0xEF, 0xF7)   # กล่องปกติ
COLOR_BOX_LINE = RGBColor(0x9C, 0xB5, 0xD1)
COLOR_ARROW    = RGBColor(0x7A, 0x91, 0xA8)
COLOR_BAD_BG   = RGBColor(0xFD, 0xEC, 0xEC)   # แถบ ❌
COLOR_BAD_LINE = RGBColor(0xD9, 0x53, 0x4F)
COLOR_GOOD_BG  = RGBColor(0xEA, 0xF6, 0xEC)   # แถบ ✅
COLOR_GOOD_LINE = RGBColor(0x3F, 0x9B, 0x54)

# กล่อง prompt — สีต่างจากบล็อกโค้ดชัดเจน เพื่อให้ผู้เรียนแยกออกทันที
# ว่า "อันนี้คือสิ่งที่ต้องพิมพ์ให้ Claude" ไม่ใช่โค้ดที่ต้องอ่าน
COLOR_PROMPT_BG   = RGBColor(0xFF, 0xF7, 0xE6)
COLOR_PROMPT_LINE = RGBColor(0xE0, 0xA0, 0x2C)
COLOR_PROMPT_LBL  = RGBColor(0xA8, 0x6C, 0x0A)

SIZE_TITLE_COVER   = Pt(40)
SIZE_TITLE_SECTION = Pt(40)
SIZE_TITLE_SLIDE   = Pt(36)
SIZE_BODY          = Pt(26)   # ห้ามต่ำกว่า 20pt
SIZE_BODY_MIN      = Pt(20)
SIZE_BIG           = Pt(54)
SIZE_CODE          = Pt(18)
SIZE_FOOTER        = Pt(12)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.9)


# ──────────────────────────── ตัวช่วยจัดฟอนต์ ────────────────────────────

def style_run(run, *, font=THAI_FONT, size=SIZE_BODY, bold=False, color=COLOR_TEXT):
    """ตั้งฟอนต์ให้ครบทั้ง latin และ complex script

    ⚠️ สำคัญมากสำหรับภาษาไทย: ถ้าไม่ตั้ง <a:cs> ด้วย สระและวรรณยุกต์จะลอยหรือซ้อนกัน
    """
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color

    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':')[1]}")
        if el is None:
            from pptx.oxml.ns import qn
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def strip_md(text):
    """ตัดสัญลักษณ์ markdown ที่ python-pptx แสดงไม่ได้ออก"""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text.strip()


# ──────────────────────────── ตัวอ่านไฟล์ต้นฉบับ ────────────────────────────

def parse(md_text):
    meta = {}
    body = md_text

    fm = re.match(r"^---\s*\n(.*?)\n---\s*\n(.*)$", md_text, re.S)
    if fm:
        for line in fm.group(1).splitlines():
            if ":" in line:
                k, v = line.split(":", 1)
                meta[k.strip()] = v.strip().strip('"').strip("'")
        body = fm.group(2)

    slides = []
    current = None
    in_code = False

    def close():
        nonlocal current
        if current:
            slides.append(current)
            current = None

    for raw in body.splitlines():
        line = raw.rstrip()

        if line.startswith("```"):
            if not in_code:
                in_code = True
                if current is None:
                    current = {"kind": "content", "title": "", "items": [], "notes": []}
                fence = line[3:].strip()
                if fence.startswith("flow"):
                    parts = fence.split()
                    shape = parts[1] if len(parts) > 1 else "chain"
                    current["items"].append({"type": "flow", "shape": shape, "lines": []})
                elif fence.startswith("prompt"):
                    label = fence[len("prompt"):].strip() or "พิมพ์แบบนี้"
                    current["items"].append({"type": "prompt", "label": label, "lines": []})
                else:
                    current["items"].append({"type": "code", "lines": []})
            else:
                in_code = False
            continue

        if in_code:
            if current and current["items"] and current["items"][-1]["type"] in ("code", "flow", "prompt"):
                current["items"][-1]["lines"].append(raw)
            continue

        if line.startswith("## "):
            close()
            slides.append({"kind": "section", "title": strip_md(line[3:]), "items": [], "notes": []})
            continue

        if line.startswith("# "):
            close()
            current = {"kind": "content", "title": strip_md(line[2:]), "items": [], "notes": []}
            continue

        if line.startswith("!! "):
            close()
            slides.append({"kind": "big", "title": strip_md(line[3:]), "items": [], "notes": []})
            continue

        if line.startswith("> note:"):
            target = current if current else (slides[-1] if slides else None)
            if target is not None:
                target["notes"].append(strip_md(line[7:]))
            continue

        img = re.match(r"!\[(.*?)\]\((.*?)\)", line.strip())
        if img:
            if current is None:
                current = {"kind": "content", "title": "", "items": [], "notes": []}
            current["items"].append({"type": "image", "caption": img.group(1), "path": img.group(2)})
            continue

        # ตาราง markdown — รวมบรรทัดที่ขึ้นต้นและลงท้ายด้วย | ให้เป็นตารางเดียว
        if line.strip().startswith("|") and line.strip().endswith("|"):
            cells = [strip_md(c) for c in line.strip().strip("|").split("|")]
            # ข้ามแถวคั่นหัวตาราง เช่น |---|:--:|
            if all(re.fullmatch(r":?-{2,}:?", c.strip()) for c in cells if c.strip()):
                continue
            if current is None:
                current = {"kind": "content", "title": "", "items": [], "notes": []}
            if current["items"] and current["items"][-1]["type"] == "table":
                current["items"][-1]["rows"].append(cells)
            else:
                current["items"].append({"type": "table", "rows": [cells]})
            continue

        if line.strip().startswith(("- ", "* ")):
            if current is None:
                current = {"kind": "content", "title": "", "items": [], "notes": []}
            indent = (len(line) - len(line.lstrip())) // 2
            current["items"].append({"type": "bullet", "text": strip_md(line.strip()[2:]), "level": min(indent, 2)})
            continue

    close()
    return meta, slides


# ──────────────────────────── ตัวสร้างสไลด์ ────────────────────────────

def build_cover(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    _, tf = add_textbox(slide, MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(2.6))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = meta.get("title", "")
    style_run(r, size=SIZE_TITLE_COVER, bold=True, color=COLOR_ACCENT)

    for key, size, color in (("subtitle", Pt(22), COLOR_TEXT), ("session_date", Pt(18), COLOR_MUTED), ("date", Pt(18), COLOR_MUTED)):
        if meta.get(key):
            p2 = tf.add_paragraph()
            p2.space_before = Pt(14)
            r = p2.add_run()
            r.text = meta[key]
            style_run(r, size=size, color=color)


def build_section(prs, s, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_ACCENT
    bg.line.fill.background()

    _, tf = add_textbox(slide, MARGIN, Inches(3.0), SLIDE_W - 2 * MARGIN, Inches(1.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s["title"]
    style_run(r, size=SIZE_TITLE_SECTION, bold=True, color=COLOR_WHITE)

    attach_notes(slide, s)


def build_big(prs, s, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ข้อความยาวให้ลดขนาดลงจนพอดี 3 บรรทัด ไม่ให้ล้นกรอบ
    width_in = Emu(SLIDE_W - 2 * MARGIN).inches - 0.3
    size = SIZE_BIG
    for pt_size in (54, 46, 40, 34, 30):
        if est_lines(s["title"], width_in, pt_size) <= 3:
            size = Pt(pt_size)
            break
        size = Pt(pt_size)

    _, tf = add_textbox(slide, MARGIN, Inches(2.2), SLIDE_W - 2 * MARGIN, Inches(3.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s["title"]
    style_run(r, size=size, bold=True, color=COLOR_ACCENT)
    attach_notes(slide, s)


def est_lines(text, width_in, size_pt, factor=0.58):
    """ประมาณว่าข้อความจะตัดเป็นกี่บรรทัดในกล่องกว้างเท่านี้

    ⚠️ ค่า factor ต้องตรงหรือมากกว่าที่ scripts/check_slides.py ใช้ (0.55)
       เพื่อให้ตัวสร้างจองพื้นที่เผื่อไว้มากกว่าที่ตัวตรวจคาดเสมอ
       ไม่งั้นจะสร้างผ่านแต่ตรวจไม่ผ่าน
    """
    char_w_in = size_pt * factor / 72
    per_line = max(1, int(width_in / char_w_in))
    return max(1, -(-len(text) // per_line))


def _bullet_height(item, size):
    """ความสูงที่ bullet ข้อนี้ต้องใช้จริง (นับกรณีตัดหลายบรรทัดด้วย)"""
    avail_in = Emu(SLIDE_W - 2 * MARGIN - Inches(0.3 * item["level"])).inches - 0.2
    n = est_lines("• " + item["text"], avail_in, size.pt)
    return Inches(size.pt * 1.42 / 72) * n + Inches(0.16)


def build_content(prs, s, index, md_dir, missing_images, overflow_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    top = MARGIN
    if s["title"]:
        # หัวข้อยาวให้ลดขนาดลงก่อน ถ้ายังยาวก็ขยายกล่องตามจำนวนบรรทัดจริง
        w_in = Emu(SLIDE_W - 2 * MARGIN).inches - 0.2
        t_size = SIZE_TITLE_SLIDE
        for pt_size in (36, 32, 28):
            t_size = Pt(pt_size)
            if est_lines(s["title"], w_in, pt_size) <= 2:
                break
        n_lines = est_lines(s["title"], w_in, t_size.pt)
        title_h = Inches(t_size.pt * 1.32 / 72) * n_lines + Inches(0.22)

        _, tf = add_textbox(slide, MARGIN, Inches(0.55), SLIDE_W - 2 * MARGIN, title_h)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = s["title"]
        style_run(r, size=t_size, bold=True, color=COLOR_ACCENT)
        top = Inches(0.55) + title_h + Inches(0.18)

    bullets = [i for i in s["items"] if i["type"] == "bullet"]

    # หาขนาดตัวอักษรที่ทำให้ bullet ทั้งหมดพอดีหน้า
    # (bullet ที่ยาวเกิน 1 บรรทัดต้องได้ความสูงเพิ่มตามจริง ไม่งั้นข้อความจะทับกัน)
    body_size = SIZE_BODY
    if bullets:
        avail_for_bullets = SLIDE_H - top - Inches(0.75)
        other = sum(1 for i in s["items"] if i["type"] in ("code", "prompt", "table", "flow", "image"))
        for pt_size in (26, 24, 22, 20):
            total = Emu(0)
            for b in bullets:
                total += _bullet_height(b, Pt(pt_size))
            if other == 0 and total <= avail_for_bullets:
                body_size = Pt(pt_size)
                break
            body_size = Pt(pt_size)
            if other and total <= avail_for_bullets * 0.55:
                break
        if other == 0 and sum((_bullet_height(b, body_size) for b in bullets), Emu(0)) > avail_for_bullets:
            overflow_slides.append(
                f"หน้า {index + 1} · {s['title'] or '(ไม่มีหัวข้อ)'} — bullet {len(bullets)} ข้อ ยาวเกินหน้า"
            )

    cursor = top
    for item in s["items"]:
        if item["type"] == "bullet":
            h = _bullet_height(item, body_size)
            _, tf = add_textbox(slide, MARGIN + Inches(0.3 * item["level"]),
                                cursor, SLIDE_W - 2 * MARGIN - Inches(0.3 * item["level"]),
                                h - Inches(0.06))
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = ("• " if item["level"] == 0 else "– ") + item["text"]
            style_run(r, size=body_size)
            cursor += h

        elif item["type"] == "code":
            lines = item["lines"] or [""]

            # ย่อให้พอดีสไลด์อัตโนมัติ — บล็อกโค้ดยาว (เช่น prompt) จะได้ไม่ล้นขอบ
            avail = SLIDE_H - cursor - Inches(0.75)
            line_h = Inches(0.34)
            code_size = SIZE_CODE
            need = line_h * len(lines) + Inches(0.4)
            if need > avail and len(lines):
                scale = max(0.55, (avail - Inches(0.3)) / (line_h * len(lines)))
                line_h = int(line_h * scale)
                code_size = Pt(max(10, round(SIZE_CODE.pt * scale)))
                need = line_h * len(lines) + Inches(0.3)
                if need > avail:
                    overflow_slides.append(
                        f"หน้า {index + 1} · {s['title'] or '(ไม่มีหัวข้อ)'} — บล็อกโค้ด {len(lines)} บรรทัด ยาวเกินหน้า"
                    )

            h = need
            box = slide.shapes.add_shape(1, MARGIN, cursor, SLIDE_W - 2 * MARGIN, h)
            box.fill.solid()
            box.fill.fore_color.rgb = COLOR_CODE_BG
            box.line.fill.background()
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = Inches(0.25)
            tf.margin_top = Inches(0.15)
            for n, ln in enumerate(lines):
                p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = ln
                style_run(r, font=MONO_FONT, size=code_size)
            cursor += h + Inches(0.2)

        elif item["type"] == "prompt":
            # กล่อง prompt — ใช้ฟอนต์ไทย ไม่ใช่ฟอนต์โค้ด เพราะเป็นประโยคภาษาไทยที่ผู้เรียนพิมพ์เอง
            lines = [ln.rstrip() for ln in item["lines"]] or [""]
            avail = SLIDE_H - cursor - Inches(0.75)
            line_h = Inches(0.42)
            p_size = Pt(20)
            head_h = Inches(0.42)
            need = head_h + line_h * len(lines) + Inches(0.34)
            if need > avail and len(lines):
                scale = max(0.6, (avail - head_h - Inches(0.34)) / (line_h * len(lines)))
                line_h = int(line_h * scale)
                p_size = Pt(max(13, round(20 * scale)))
                need = head_h + line_h * len(lines) + Inches(0.3)
                if need > avail:
                    overflow_slides.append(
                        f"หน้า {index + 1} · {s['title'] or '(ไม่มีหัวข้อ)'} — กล่อง prompt {len(lines)} บรรทัด ยาวเกินหน้า"
                    )

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         MARGIN, cursor, SLIDE_W - 2 * MARGIN, need)
            box.fill.solid()
            box.fill.fore_color.rgb = COLOR_PROMPT_BG
            box.line.color.rgb = COLOR_PROMPT_LINE
            box.line.width = Pt(1.75)
            box.shadow.inherit = False
            box.text_frame.text = ""

            # แถบสีด้านซ้าย ให้สะดุดตาแยกจากบล็อกโค้ด
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         MARGIN, cursor, Inches(0.11), need)
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_PROMPT_LINE
            bar.line.fill.background()
            bar.shadow.inherit = False

            # ป้ายหัวกล่อง
            _, tf = add_textbox(slide, MARGIN + Inches(0.3), cursor + Inches(0.06),
                                SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.34))
            p = tf.paragraphs[0]
            r = p.add_run()
            p.alignment = PP_ALIGN.LEFT
            r.text = f"💬 {item['label']}"
            style_run(r, size=Pt(14), bold=True, color=COLOR_PROMPT_LBL)

            # เนื้อ prompt
            _, tf = add_textbox(slide, MARGIN + Inches(0.3), cursor + head_h,
                                SLIDE_W - 2 * MARGIN - Inches(0.5),
                                line_h * len(lines))
            tf.word_wrap = True
            for n, ln in enumerate(lines):
                p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
                p.space_after = Pt(0)
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = ln
                style_run(r, size=p_size, color=COLOR_TEXT)

            cursor += need + Inches(0.22)

        elif item["type"] == "flow":
            avail = SLIDE_H - cursor - Inches(0.75)
            used = draw_flow(slide, item, MARGIN, cursor + Inches(0.1),
                             SLIDE_W - 2 * MARGIN, avail)
            cursor += used + Inches(0.35)

        elif item["type"] == "table":
            rows = item["rows"]
            n_rows = len(rows)
            n_cols = max(len(r) for r in rows)
            avail_h = SLIDE_H - cursor - Inches(0.8)
            row_h = min(Inches(0.55), max(Inches(0.32), avail_h / n_rows))
            h = row_h * n_rows

            shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, cursor,
                                           SLIDE_W - 2 * MARGIN, h)
            tbl = shape.table
            tbl.first_row = True

            # ขนาดตัวอักษรลดตามจำนวนแถว แต่ไม่ต่ำกว่า 12pt เพื่อให้ยังอ่านออกบนจอ
            t_size = Pt(18) if n_rows <= 6 else (Pt(15) if n_rows <= 9 else Pt(12))

            for r_i, row in enumerate(rows):
                tbl.rows[r_i].height = row_h
                for c_i in range(n_cols):
                    cell = tbl.cell(r_i, c_i)
                    cell.margin_left = Inches(0.1)
                    cell.margin_right = Inches(0.1)
                    cell.margin_top = Inches(0.03)
                    cell.margin_bottom = Inches(0.03)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_ACCENT if r_i == 0 else (
                        COLOR_WHITE if r_i % 2 else COLOR_CODE_BG)

                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = row[c_i] if c_i < len(row) else ""
                    style_run(run, size=t_size, bold=(r_i == 0),
                              color=COLOR_WHITE if r_i == 0 else COLOR_TEXT)

            cursor += h + Inches(0.2)

        elif item["type"] == "image":
            img_path = (md_dir / item["path"]).resolve()
            h = Inches(3.2)
            if img_path.exists():
                slide.shapes.add_picture(str(img_path), MARGIN, cursor, height=h)
            else:
                missing_images.append(f"{item['path']}  —  {item['caption']}")
                box = slide.shapes.add_shape(1, MARGIN, cursor, SLIDE_W - 2 * MARGIN, h)
                box.fill.solid()
                box.fill.fore_color.rgb = COLOR_IMG_BG
                box.line.fill.background()
                tf = box.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = f"[รอภาพ]  {item['caption']}"
                style_run(r, size=Pt(18), color=COLOR_MUTED)
            cursor += h + Inches(0.25)

    attach_notes(slide, s)


# ──────────────────────────── รูปวาด (diagram) ────────────────────────────
#
# รูปวาดทั้งหมดเป็น "รูปทรงจริงของ PowerPoint" ไม่ใช่ไฟล์ภาพ
# จึงคมทุกระดับซูม · ผู้สอนแก้ข้อความในสไลด์ได้เอง · ฟอนต์ไทยแสดงถูกต้อง
#
# รูปแบบที่รองรับ (เขียนในต้นฉบับ .md)
#
#   ```flow chain          กล่องเรียงแนวนอน มีลูกศรเชื่อม
#   ผู้ใช้กดปุ่ม | หน้าเว็บ | *Firestore | แสดงผลกลับ
#   ```
#
#   ```flow cycle          กล่องเรียงแนวนอน + ลูกศรวนกลับด้านล่าง
#   วางแผน | ลงมือ | ตรวจด้วยตา | commit
#   ```
#
#   ```flow compare        สองแถบเทียบกัน ❌ กับ ✅
#   ❌ ใส่คีย์ไว้ในหน้าเว็บ :: หน้าเว็บ *มีคีย์ | OpenRouter
#   ✅ ให้ตัวกลางถือคีย์ :: หน้าเว็บ | ตัวกลาง *มีคีย์ | OpenRouter
#   ```
#
# ใส่ * นำหน้าข้อความในกล่อง เพื่อเน้นกล่องนั้น (พื้นเข้ม ตัวอักษรขาว)


def _box(slide, left, top, w, h, text, *, highlight=False, size=Pt(14),
         fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = (COLOR_ACCENT if highlight else (fill or COLOR_BOX_BG))
    shp.line.color.rgb = line or COLOR_BOX_LINE
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    style_run(r, size=size, bold=highlight,
              color=COLOR_WHITE if highlight else COLOR_TEXT)
    return shp


def _arrow(slide, left, top, w, h):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLOR_ARROW
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _parse_items(text):
    """แยกข้อความในบรรทัดเป็นรายการกล่อง

    มี '*' อยู่ที่ไหนก็ได้ในข้อความ = กล่องนั้นเป็นกล่องเน้น (พื้นเข้ม ตัวอักษรขาว)
    แล้วตัด '*' ทิ้งก่อนแสดงผล
    """
    out = []
    for part in text.split("|"):
        t = strip_md(part.strip())
        hi = "*" in t
        t = re.sub(r"\s*\*\s*", " ", t).strip()
        if t:
            out.append((t, hi))
    return out


def _draw_chain(slide, items, left, top, width, box_h, size=Pt(14)):
    """กล่องเรียงแนวนอนพร้อมลูกศร คืนค่าความสูงที่ใช้ไป"""
    n = len(items)
    if n == 0:
        return Emu(0)
    gap = Inches(0.42) if n > 1 else Emu(0)
    box_w = int((width - gap * (n - 1)) / n)
    x = left
    for i, (text, hi) in enumerate(items):
        _box(slide, x, top, box_w, box_h, text, highlight=hi, size=size)
        x += box_w
        if i < n - 1:
            a_h = int(box_h * 0.30)
            _arrow(slide, x + int(gap * 0.12), top + int((box_h - a_h) / 2),
                   int(gap * 0.76), a_h)
            x += gap
    return box_h


def draw_flow(slide, item, left, top, width, avail_h):
    shape = item.get("shape", "chain")
    lines = [ln for ln in item["lines"] if ln.strip()]
    if not lines:
        return Emu(0)

    if shape == "compare":
        used = Emu(0)
        panel_gap = Inches(0.22)
        n_panels = len(lines)
        panel_h = min(Inches(1.65),
                      int((avail_h - panel_gap * (n_panels - 1)) / max(n_panels, 1)))
        for ln in lines:
            head, _, body = ln.partition("::")
            head = strip_md(head.strip())
            good = head.startswith("✅")
            bg   = COLOR_GOOD_BG if good else COLOR_BAD_BG
            edge = COLOR_GOOD_LINE if good else COLOR_BAD_LINE

            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           left, top + used, width, panel_h)
            panel.fill.solid()
            panel.fill.fore_color.rgb = bg
            panel.line.color.rgb = edge
            panel.line.width = Pt(1.5)
            panel.shadow.inherit = False
            panel.text_frame.text = ""

            # หัวแถบ
            _, tf = add_textbox(slide, left + Inches(0.18), top + used + Inches(0.08),
                                width - Inches(0.36), Inches(0.34))
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = head
            style_run(r, size=Pt(15), bold=True, color=edge)

            # กล่องข้างใน
            items = _parse_items(body)
            inner_h = panel_h - Inches(0.56)
            _draw_chain(slide, items,
                        left + Inches(0.28), top + used + Inches(0.46),
                        width - Inches(0.56), inner_h, size=Pt(13))

            used += panel_h + panel_gap
        return used - panel_gap

    # chain / cycle
    items = _parse_items(lines[0])
    box_h = min(Inches(1.15), avail_h - (Inches(0.6) if shape == "cycle" else Emu(0)))
    _draw_chain(slide, items, left, top, width, box_h)

    if shape == "cycle":
        bar_y = top + box_h + Inches(0.16)
        bar = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, left, bar_y,
                                     width, Inches(0.32))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ARROW
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = bar.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "วนกลับไปทำรอบต่อไป"
        style_run(r, size=Pt(12), bold=True, color=COLOR_WHITE)
        return box_h + Inches(0.48)

    return box_h


def attach_notes(slide, s):
    text = "\n".join(s["notes"]) if s["notes"] else ""
    slide.notes_slide.notes_text_frame.text = text


def add_footers(prs, meta):
    label = f"RAISE 2 · เดือนที่ 2 · สัปดาห์ที่ {meta.get('week', '')}"
    for n, slide in enumerate(prs.slides, start=1):
        if n == 1:
            continue
        _, tf = add_textbox(slide, MARGIN, SLIDE_H - Inches(0.55),
                            SLIDE_W - 2 * MARGIN, Inches(0.35))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}          {n}"
        style_run(r, size=SIZE_FOOTER, color=COLOR_MUTED)


# ──────────────────────────── main ────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("วิธีใช้: python scripts/build_slides.py <ไฟล์ .md>")
        sys.exit(1)

    md_path = Path(sys.argv[1]).resolve()
    if not md_path.exists():
        print(f"ไม่พบไฟล์: {md_path}")
        sys.exit(1)

    meta, slides = parse(md_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs, meta)

    missing_images = []
    overflow_slides = []
    for i, s in enumerate(slides, start=1):
        if s["kind"] == "section":
            build_section(prs, s, i)
        elif s["kind"] == "big":
            build_big(prs, s, i)
        else:
            build_content(prs, s, i, md_path.parent, missing_images, overflow_slides)

    add_footers(prs, meta)

    out = md_path.with_suffix(".pptx")
    prs.save(str(out))

    print(f"✅ สร้างแล้ว: {out}")
    print(f"   จำนวนสไลด์: {len(prs.slides._sldIdLst)}")

    no_notes = [i for i, s in enumerate(slides, start=2) if not s["notes"]]
    if no_notes:
        print(f"\n⚠️  สไลด์ที่ยังไม่มีบันทึกผู้บรรยาย (> note:) หน้า: {no_notes}")

    if overflow_slides:
        print(f"\n⚠️  สไลด์ที่เนื้อหายาวเกินหน้า ({len(overflow_slides)}):")
        for o in overflow_slides:
            print(f"   - {o}")
        print("   → ควรแตกเป็นสองสไลด์")

    if missing_images:
        print(f"\n📸 ภาพที่ยังต้องถ่ายเพิ่ม ({len(missing_images)} ภาพ):")
        for m in missing_images:
            print(f"   - {m}")


if __name__ == "__main__":
    main()
